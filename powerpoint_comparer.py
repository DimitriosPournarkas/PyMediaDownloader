import sys
import os
from pptx import Presentation
import re

def compare_pptx_files(file1, file2):
    try:
        # Decode file paths for special characters
        file1 = file1.encode('latin-1').decode('utf-8')
        file2 = file2.encode('latin-1').decode('utf-8')
        
        
        # Read PowerPoint presentations
        prs1 = Presentation(file1)
        prs2 = Presentation(file2)
        
        # Extract text from both presentations
        text1 = extract_text(prs1)
        text2 = extract_text(prs2)
        
        # Calculate similarity
        similarity = calculate_text_similarity(text1, text2)
        
        return similarity > 0.6
        
    except Exception as e:
        return False

def extract_text(prs):
    """Extract text from PowerPoint presentation"""
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def calculate_text_similarity(text1, text2):
    """Calculate similarity between two texts"""
    if not text1 or not text2:
        return 0.0
    
    words1 = set(re.findall(r'\w+', text1.lower()))
    words2 = set(re.findall(r'\w+', text2.lower()))
    
    if not words1 or not words2:
        return 0.0
    
    common_words = words1.intersection(words2)
    similarity = len(common_words) / max(len(words1), len(words2))
    
    return similarity

if __name__ == "__main__":
    if len(sys.argv) != 3:
        
        sys.exit(1)
    
    file1, file2 = sys.argv[1], sys.argv[2]
    similar = compare_pptx_files(file1, file2)
    sys.exit(0 if similar else 1)